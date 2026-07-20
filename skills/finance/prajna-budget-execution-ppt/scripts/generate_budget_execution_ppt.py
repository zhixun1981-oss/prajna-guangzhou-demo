#!/usr/bin/env python3
"""
prajna-budget-execution-ppt 生成器 v1.0.0
为管理层一键生成「本月预算执行情况汇报」PPTX。
包含：预算完成率对比图、超支项目分析、结余资金使用建议、下月预算调整方案。
"""

import argparse
import json
import os
import re
import sys
from datetime import datetime
from pathlib import Path

# Graceful dependency handling
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    HAVE_PPTX = True
except ImportError as exc:
    HAVE_PPTX = False
    PPTX_IMPORT_ERROR = str(exc)


# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
SCRIPT_DIR = Path(__file__).resolve().parent
SKILL_DIR = SCRIPT_DIR.parent
DATA_DIR = SKILL_DIR / "data"
SAMPLES_DIR = Path.home() / ".prajna" / "prajna-budget-execution-ppt" / "samples"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
CN_FONT = "Microsoft YaHei"
ACCENT_BLUE = RGBColor(0x1F, 0x4E, 0x78)
DARK_BLUE = RGBColor(0x12, 0x2A, 0x4D)
LIGHT_BLUE = RGBColor(0x44, 0x72, 0xC4)
GOLD = RGBColor(0xC5, 0xA0, 0x4E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x88, 0x88, 0x88)
RED = RGBColor(0xC0, 0x00, 0x00)
GREEN = RGBColor(0x00, 0x80, 0x00)
ORANGE = RGBColor(0xFF, 0x66, 0x00)


def load_json(path):
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def _default_data():
    """Fallback data when the sample JSON is missing."""
    return {
        "month": "2026年7月",
        "company": "启明星科技有限公司",
        "author": "财务部",
        "currency": "CNY",
        "categories": [
            {"name": "人力成本", "budget": 1200000, "actual": 1250000, "dept": "人力资源部/财务部", "reason": "招聘进度超预期、加班费增加及社保公积金基数上调", "severity": "high"},
            {"name": "市场推广费", "budget": 300000, "actual": 345000, "dept": "市场部", "reason": "行业展会与线上投放集中上线", "severity": "high"},
            {"name": "差旅费", "budget": 180000, "actual": 165000, "dept": "销售部/各部门", "reason": "部分客户拜访改为线上会议", "severity": "normal"},
            {"name": "办公费", "budget": 80000, "actual": 62000, "dept": "行政部", "reason": "推行无纸化办公及耗材集采", "severity": "normal"},
            {"name": "设备采购", "budget": 250000, "actual": 220000, "dept": "IT部/行政部", "reason": "部分设备采购延期至下月验收", "severity": "normal"},
            {"name": "云服务与IT支出", "budget": 150000, "actual": 158000, "dept": "IT部", "reason": "业务量增长导致服务器扩容", "severity": "medium"},
            {"name": "培训费", "budget": 60000, "actual": 35000, "dept": "人力资源部", "reason": "外部讲师档期推迟", "severity": "normal"},
            {"name": "业务招待费", "budget": 50000, "actual": 58000, "dept": "销售部", "reason": "重点客户接待频次增加", "severity": "medium"},
        ],
        "root_causes": [
            {"category": "外部市场因素", "detail": "行业展会档期集中、流量成本上涨，导致市场推广费前置。"},
            {"category": "业务增长驱动", "detail": "新业务线扩招及云资源扩容，带动人力与IT支出同步增长。"},
            {"category": "费用集中投放", "detail": "品牌活动与重点客户接待集中在7月执行，形成阶段性超支。"},
            {"category": "内部管控疏漏", "detail": "部分设备采购审批周期较长，培训计划执行跟踪不足。"},
        ],
        "next_month_adjustments": [
            {"name": "市场推广费", "original_budget": 300000, "adjustment": 35000, "reason": "延续展会及线上投放需求"},
            {"name": "云服务与IT支出", "original_budget": 150000, "adjustment": 15000, "reason": "预留服务器扩容及安全加固费用"},
            {"name": "业务招待费", "original_budget": 50000, "adjustment": 8000, "reason": "配合Q3大客户攻坚"},
            {"name": "培训费", "original_budget": 60000, "adjustment": -10000, "reason": "部分课程改为内部讲师，压缩外采预算"},
            {"name": "办公费", "original_budget": 80000, "adjustment": -5000, "reason": "无纸化办公持续降本"},
        ],
        "surplus_suggestions": [
            {"name": "差旅费", "surplus": 15000, "usage": "补充重点客户现场拜访基金", "benefit": "提升大客户转化率", "dept": "销售部"},
            {"name": "办公费", "surplus": 18000, "usage": "升级会议室智能设备", "benefit": "提升远程会议体验", "dept": "行政部"},
            {"name": "设备采购", "surplus": 30000, "usage": "提前采购下季度研发测试机", "benefit": "缩短研发设备到位周期", "dept": "IT部"},
            {"name": "培训费", "surplus": 25000, "usage": "投入AI工具内部训练营", "benefit": "提升人均产出约8%", "dept": "人力资源部"},
        ],
    }


def load_data(data_path=None):
    if data_path:
        return load_json(Path(data_path))
    default_path = DATA_DIR / "budget_execution_sample.json"
    if default_path.exists():
        return load_json(default_path)
    return _default_data()


def _safe_filename(name):
    name = re.sub(r"[\\/:*?\"<>|\s]+", "_", str(name))
    return name.strip("_") or "report"


def _set_text(text_frame, text, font_name=CN_FONT, font_size=18, bold=False,
              color=BLACK, align=PP_ALIGN.LEFT):
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name=CN_FONT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    shape.line.fill.background()
    _set_text(shape.text_frame, text, font_name=font_name, font_size=font_size,
              bold=bold, color=color, align=align)
    return shape


def _add_title_bar(slide, title_text, subtitle_text=""):
    """Top dark-blue title bar used on content slides."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                  Inches(13.333), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()
    _set_text(bar.text_frame, title_text, font_size=24, bold=True,
              color=WHITE, align=PP_ALIGN.LEFT)
    bar.text_frame.margin_left = Inches(0.5)
    bar.text_frame.margin_top = Inches(0.25)
    if subtitle_text:
        sub = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.7),
                                      Inches(12), Inches(0.35))
        sub.fill.background()
        sub.line.fill.background()
        _set_text(sub.text_frame, subtitle_text, font_size=12, color=WHITE,
                  align=PP_ALIGN.LEFT)


def _fmt_money(value):
    return f"{value:,.0f}"


def _fmt_pct(value):
    return f"{value * 100:.1f}%"


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def build_cover(prs, data):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE

    # Decorative gold top line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.7),
                                   Inches(13.333), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    # Title
    title_box = _add_text_box(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.5),
                              "本月预算执行情况汇报", font_size=54, bold=True,
                              color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    _add_text_box(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.6),
                  f"{data['company']}  ·  {data['month']}", font_size=24,
                  color=GOLD, align=PP_ALIGN.CENTER)

    # Meta
    meta_text = f"汇报部门：{data.get('author', '财务部')}\n汇报日期：{datetime.now().strftime('%Y-%m-%d')}"
    _add_text_box(slide, Inches(4.5), Inches(5.5), Inches(4.3), Inches(1),
                  meta_text, font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

    # Bottom accent shape
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8),
                                  Inches(13.333), Inches(0.7))
    bot.fill.solid()
    bot.fill.fore_color.rgb = LIGHT_BLUE
    bot.line.fill.background()
    _set_text(bot.text_frame, "PRAJNA 企智 · 财务预算智能辅助", font_size=14,
              color=WHITE, align=PP_ALIGN.CENTER)
    return slide


def build_summary(slide, data):
    _add_title_bar(slide, "本月预算执行总览", "总体进度、偏差与结构一目了然")

    cats = data["categories"]
    total_budget = sum(c["budget"] for c in cats)
    total_actual = sum(c["actual"] for c in cats)
    variance = total_actual - total_budget
    completion = total_actual / total_budget if total_budget else 0
    overspend_count = sum(1 for c in cats if c["actual"] > c["budget"])
    surplus_count = sum(1 for c in cats if c["actual"] < c["budget"])

    metrics = [
        ("月度总预算", f"{total_budget:,.0f} 元", "预算批复总额"),
        ("实际支出", f"{total_actual:,.0f} 元", "截至本月末实际入账"),
        ("总体完成率", f"{completion * 100:.1f}%", "实际/预算"),
        ("总体偏差额", f"{variance:+,.0f} 元", "超支为正，结余为负"),
        ("超支项目数", f"{overspend_count} 项", "需重点分析原因"),
        ("结余项目数", f"{surplus_count} 项", "可统筹再分配"),
    ]

    left_start = Inches(0.8)
    top_start = Inches(1.6)
    box_w = Inches(3.8)
    box_h = Inches(1.4)
    gap_x = Inches(0.4)
    gap_y = Inches(0.35)

    for idx, (label, value, note) in enumerate(metrics):
        row = idx // 3
        col = idx % 3
        left = left_start + col * (box_w + gap_x)
        top = top_start + row * (box_h + gap_y)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF3, 0xF6, 0xFA)
        box.line.color.rgb = ACCENT_BLUE
        box.line.width = Pt(1)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_bottom = Inches(0.15)

        p1 = tf.paragraphs[0]
        p1.text = label
        p1.alignment = PP_ALIGN.LEFT
        run1 = p1.runs[0]
        run1.font.name = CN_FONT
        run1.font.size = Pt(14)
        run1.font.color.rgb = GRAY

        p2 = tf.add_paragraph()
        p2.text = value
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(8)
        run2 = p2.runs[0]
        run2.font.name = CN_FONT
        run2.font.size = Pt(26)
        run2.font.bold = True
        run2.font.color.rgb = ACCENT_BLUE if variance >= 0 else GREEN

        p3 = tf.add_paragraph()
        p3.text = note
        p3.alignment = PP_ALIGN.LEFT
        p3.space_before = Pt(4)
        run3 = p3.runs[0]
        run3.font.name = CN_FONT
        run3.font.size = Pt(11)
        run3.font.color.rgb = GRAY


def build_budget_actual_chart(slide, data):
    _add_title_bar(slide, "预算完成率对比", "各费用项目预算 vs 实际支出")

    cats = data["categories"]
    names = [c["name"] for c in cats]
    budgets = [c["budget"] / 10000 for c in cats]
    actuals = [c["actual"] / 10000 for c in cats]

    chart_data = CategoryChartData()
    chart_data.categories = names
    chart_data.add_series("预算（万元）", budgets)
    chart_data.add_series("实际（万元）", actuals)

    x, y, cx, cy = Inches(0.7), Inches(1.5), Inches(12), Inches(5.3)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = "费用项目预算执行对比"
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM


def build_completion_rate_chart(slide, data):
    _add_title_bar(slide, "完成率分布与排名", "100% 基准线辅助识别超支/结余")

    cats = sorted(data["categories"], key=lambda c: c["actual"] / c["budget"], reverse=True)
    names = [c["name"] for c in cats]
    rates = [(c["actual"] / c["budget"]) * 100 for c in cats]
    threshold = [100] * len(names)

    chart_data = CategoryChartData()
    chart_data.categories = names
    chart_data.add_series("完成率(%)", rates)
    chart_data.add_series("预算基准线", threshold)

    x, y, cx, cy = Inches(0.7), Inches(1.5), Inches(12), Inches(5.3)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = "费用完成率排名（%）"
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM


def build_overspend_table(slide, data):
    _add_title_bar(slide, "超支项目分析", "聚焦偏差金额与责任部门")

    cats = [c for c in data["categories"] if c["actual"] > c["budget"]]
    cats.sort(key=lambda c: c["actual"] - c["budget"], reverse=True)

    rows = len(cats) + 1
    cols = 7
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.3)
    height = Inches(0.55 * rows)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    headers = ["费用项目", "预算(元)", "实际(元)", "偏差额(元)", "完成率", "责任部门", "偏差原因"]

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        _style_header_cell(cell)

    col_widths = [Inches(1.5), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.0), Inches(1.8), Inches(3.8)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for r, c in enumerate(cats, start=1):
        variance = c["actual"] - c["budget"]
        rate = c["actual"] / c["budget"]
        row_vals = [
            c["name"],
            _fmt_money(c["budget"]),
            _fmt_money(c["actual"]),
            f"+{_fmt_money(variance)}",
            _fmt_pct(rate),
            c["dept"],
            c.get("reason", "-"),
        ]
        for col_idx, val in enumerate(row_vals):
            cell = table.cell(r, col_idx)
            cell.text = val
            _style_data_cell(cell, is_warn=True)


def build_root_cause(slide, data):
    _add_title_bar(slide, "偏差根因分析", "四类因素驱动本月预算偏差")

    causes = data.get("root_causes", [])
    if not causes:
        return

    colors = [ACCENT_BLUE, LIGHT_BLUE, GOLD, RGBColor(0x70, 0xAD, 0x47)]
    left_start = Inches(0.8)
    top = Inches(1.7)
    box_w = Inches(5.8)
    box_h = Inches(2.0)
    gap_x = Inches(0.4)
    gap_y = Inches(0.4)

    for idx, cause in enumerate(causes):
        row = idx // 2
        col = idx % 2
        left = left_start + col * (box_w + gap_x)
        top_pos = top + row * (box_h + gap_y)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top_pos, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = colors[idx % len(colors)]
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)
        tf.margin_right = Inches(0.2)

        p1 = tf.paragraphs[0]
        p1.text = cause["category"]
        p1.alignment = PP_ALIGN.LEFT
        run1 = p1.runs[0]
        run1.font.name = CN_FONT
        run1.font.size = Pt(18)
        run1.font.bold = True
        run1.font.color.rgb = WHITE

        p2 = tf.add_paragraph()
        p2.text = cause["detail"]
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(8)
        run2 = p2.runs[0]
        run2.font.name = CN_FONT
        run2.font.size = Pt(14)
        run2.font.color.rgb = WHITE


def build_surplus_table(slide, data):
    _add_title_bar(slide, "结余资金使用建议", "将沉淀资金导向高价值业务场景")

    suggestions = data.get("surplus_suggestions", [])
    if not suggestions:
        # auto-generate from categories
        for c in data["categories"]:
            if c["actual"] < c["budget"]:
                suggestions.append({
                    "name": c["name"],
                    "surplus": c["budget"] - c["actual"],
                    "usage": "建议补充核心项目投入",
                    "benefit": "提升资源使用效率",
                    "dept": c["dept"],
                })

    rows = len(suggestions) + 1
    cols = 5
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.3)
    height = Inches(0.6 * rows)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    headers = ["费用项目", "结余金额(元)", "建议使用方向", "预期收益", "责任部门"]

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        _style_header_cell(cell)

    col_widths = [Inches(1.8), Inches(1.6), Inches(4.0), Inches(3.0), Inches(1.9)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for r, s in enumerate(suggestions, start=1):
        row_vals = [
            s["name"],
            _fmt_money(s["surplus"]),
            s.get("usage", ""),
            s.get("benefit", ""),
            s.get("dept", ""),
        ]
        for col_idx, val in enumerate(row_vals):
            cell = table.cell(r, col_idx)
            cell.text = val
            _style_data_cell(cell, is_warn=False)


def build_adjustment_table(slide, data):
    _add_title_bar(slide, "下月预算调整方案", "基于本月执行偏差的滚动预算修正")

    adjustments = data.get("next_month_adjustments", [])
    if not adjustments:
        return

    rows = len(adjustments) + 2  # header + data rows + total row
    cols = 5
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.3)
    height = Inches(0.6 * rows)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    headers = ["费用项目", "原预算(元)", "调整金额(元)", "调整后预算(元)", "调整说明"]

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        _style_header_cell(cell)

    col_widths = [Inches(2.0), Inches(1.6), Inches(1.6), Inches(1.8), Inches(5.3)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    total_original = 0
    total_adjusted = 0
    for r, adj in enumerate(adjustments, start=1):
        original = adj["original_budget"]
        delta = adj["adjustment"]
        adjusted = original + delta
        total_original += original
        total_adjusted += adjusted
        row_vals = [
            adj["name"],
            _fmt_money(original),
            f"{delta:+,.0f}",
            _fmt_money(adjusted),
            adj.get("reason", ""),
        ]
        for col_idx, val in enumerate(row_vals):
            cell = table.cell(r, col_idx)
            cell.text = val
            _style_data_cell(cell, is_warn=(delta > 0))

    # Total row
    total_r = rows - 1
    total_cells = ["合计", _fmt_money(total_original), _fmt_money(total_adjusted - total_original), _fmt_money(total_adjusted), ""]
    for col_idx, val in enumerate(total_cells):
        cell = table.cell(total_r, col_idx)
        cell.text = val
        _style_data_cell(cell, is_total=True)


def build_closing(slide, data):
    _add_title_bar(slide, "总结与下一步行动", "")

    cats = data["categories"]
    total_budget = sum(c["budget"] for c in cats)
    total_actual = sum(c["actual"] for c in cats)
    variance = total_actual - total_budget
    completion = total_actual / total_budget if total_budget else 0

    bullets = [
        f"核心结论：本月整体完成率 {completion * 100:.1f}%，偏差额 {_fmt_money(variance)} 元，处于可控范围。",
        "超支重点：人力成本与市场推广费为主要超支点，需复盘招聘节奏与投放ROI。",
        "结余机会：差旅、培训、设备采购形成资金沉淀，建议定向投入大客户与数字化能力建设。",
        "下月重点：滚动调整预算、强化费用事前审批、建立周度预算预警机制。",
    ]

    top = Inches(1.7)
    for idx, text in enumerate(bullets):
        left = Inches(0.8)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left,
                                      top + idx * Inches(1.1), Inches(12), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF3, 0xF6, 0xFA)
        box.line.color.rgb = ACCENT_BLUE
        box.line.width = Pt(1)
        _set_text(box.text_frame, text, font_size=16, color=BLACK, align=PP_ALIGN.LEFT)
        box.text_frame.margin_left = Inches(0.2)
        box.text_frame.margin_top = Inches(0.2)

    _add_text_box(slide, Inches(0), Inches(6.4), Inches(13.333), Inches(0.8),
                  "感谢聆听 · 欢迎交流", font_size=28, bold=True,
                  color=ACCENT_BLUE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Cell styling
# ---------------------------------------------------------------------------
def _style_header_cell(cell):
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_BLUE
    tf = cell.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].runs[0]
    run.font.name = CN_FONT
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = WHITE


def _style_data_cell(cell, is_warn=False, is_total=False):
    if is_total:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xE7, 0xE6, 0xE6)
    tf = cell.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if not tf.paragraphs[0].runs:
        tf.paragraphs[0].add_run()
    run = tf.paragraphs[0].runs[0]
    run.font.name = CN_FONT
    run.font.size = Pt(11)
    run.font.bold = is_total
    if is_warn:
        run.font.color.rgb = RED
    else:
        run.font.color.rgb = BLACK


# ---------------------------------------------------------------------------
# Main build
# ---------------------------------------------------------------------------
def build_presentation(data):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_cover(prs, data)

    summary_slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_summary(summary_slide, data)

    chart_slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    build_budget_actual_chart(chart_slide1, data)

    chart_slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    build_completion_rate_chart(chart_slide2, data)

    over_slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_overspend_table(over_slide, data)

    root_slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_root_cause(root_slide, data)

    surplus_slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_surplus_table(surplus_slide, data)

    adjust_slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_adjustment_table(adjust_slide, data)

    closing_slide = prs.slides.add_slide(prs.slide_layouts[6])
    build_closing(closing_slide, data)

    return prs


def merge_cli_args(data, args):
    if args.month:
        data["month"] = args.month
    if args.company:
        data["company"] = args.company
    if args.author:
        data["author"] = args.author
    return data


def parse_args(argv=None):
    parser = argparse.ArgumentParser(
        description="生成「本月预算执行情况汇报」PPTX"
    )
    parser.add_argument(
        "--month", "-m",
        help="汇报月份，例如 '2026年7月'，默认使用数据文件中的月份或当前月份"
    )
    parser.add_argument(
        "--company", "-c",
        help="公司名称，覆盖默认值"
    )
    parser.add_argument(
        "--author", "-a",
        help="汇报部门/作者"
    )
    parser.add_argument(
        "--data", "-d",
        help="自定义数据 JSON 文件路径，结构参考 data/budget_execution_sample.json"
    )
    parser.add_argument(
        "--output", "-o",
        help="输出文件完整路径（覆盖默认 samples 目录）"
    )
    return parser.parse_args(argv)


def main(argv=None):
    if not HAVE_PPTX:
        print(
            "[跳过] 生成 PPTX 需要 python-pptx，但当前环境未安装。\n"
            "请运行：pip install python-pptx\n"
            f"原始错误：{PPTX_IMPORT_ERROR}",
            file=sys.stderr,
        )
        return 0

    args = parse_args(argv)
    data = load_data(args.data)
    data = merge_cli_args(data, args)

    prs = build_presentation(data)

    if args.output:
        output_path = Path(args.output)
    else:
        SAMPLES_DIR.mkdir(parents=True, exist_ok=True)
        safe_company = _safe_filename(data["company"])
        safe_month = _safe_filename(data["month"])
        filename = f"本月预算执行情况汇报_{safe_company}_{safe_month}.pptx"
        output_path = SAMPLES_DIR / filename

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"已生成预算执行汇报 PPTX：{output_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
